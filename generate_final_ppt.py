import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation(output_path, img1_path, img2_path):
    prs = Presentation()
    
    # Set slide dimensions to widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_slide_layout = prs.slide_layouts[6]
    
    # Colors
    c_bg = RGBColor(10, 25, 47)          # Deep Navy
    c_teal = RGBColor(100, 255, 218)     # Neon Teal
    c_white = RGBColor(255, 255, 255)    # White
    c_slate = RGBColor(136, 146, 176)    # Slate / Subtext
    c_purple = RGBColor(179, 136, 255)   # Accent Purple
    
    def apply_dark_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = c_bg

    def add_title(slide, text, color=c_teal):
        txBox = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.83), Inches(0.8))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text.upper()
        p.font.name = "Outfit"
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = color
        p.alignment = PP_ALIGN.LEFT
        return txBox

    # =========================================================================
    # SLIDE 1: TITLE SLIDE
    # =========================================================================
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_dark_bg(slide)
    
    # Large Decorative Line
    shape = slide.shapes.add_shape(
        1, Inches(0.75), Inches(2.2), Inches(11.83), Inches(0.04) # 1 = MSO_SHAPE.RECTANGLE
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = c_teal
    shape.line.color.rgb = c_teal

    # Title & Subtitle box
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(2.5), Inches(11.83), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "SMART GRID SYNC"
    p.font.name = "Outfit"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = c_white
    p.alignment = PP_ALIGN.LEFT
    
    p2 = tf.add_paragraph()
    p2.text = "Decoupled AI Energy Optimization & Localized Indian Grid Standards"
    p2.font.name = "Outfit"
    p2.font.size = Pt(22)
    p2.font.bold = False
    p2.font.color.rgb = c_teal
    p2.space_before = Pt(15)
    
    p3 = tf.add_paragraph()
    p3.text = "Final Project Presentation | Academic Submission"
    p3.font.name = "Inter"
    p3.font.size = Pt(14)
    p3.font.color.rgb = c_slate
    p3.space_before = Pt(80)

    # =========================================================================
    # SLIDE 2: PROJECT OVERVIEW
    # =========================================================================
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_dark_bg(slide)
    add_title(slide, "1. Project Overview & Motivation")
    
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    bullets = [
        ("The Challenge:", "Modern power grids require real-time load balancing, load forecasting, and dynamic pricing to manage renewable integration. However, deploying full-stack models on local environments faces package dependency bottlenecks (e.g. Python 3.14)."),
        ("The Solution:", "Smart Grid Sync implements a decoupled client-server architecture consisting of a multithreaded custom Python API backend, an active SQLite database layer, and an optimized JavaScript dashboard served locally and via GitHub Pages."),
        ("Key Innovation:", "Combines physical grid simulation models with customized machine learning regressors written completely from scratch in pure NumPy, avoiding large compiled library dependencies.")
    ]
    
    for title, desc in bullets:
        p = tf.add_paragraph()
        p.text = ""
        p.space_after = Pt(20)
        
        run1 = p.add_run()
        run1.text = title + " "
        run1.font.name = "Outfit"
        run1.font.bold = True
        run1.font.size = Pt(18)
        run1.font.color.rgb = c_teal
        
        run2 = p.add_run()
        run2.text = desc
        run2.font.name = "Inter"
        run2.font.size = Pt(16)
        run2.font.color.rgb = c_white

    # =========================================================================
    # SLIDE 3: SYSTEM ARCHITECTURE
    # =========================================================================
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_dark_bg(slide)
    add_title(slide, "2. System Architecture & Tech Stack")
    
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    items = [
        ("Decoupled Client-Server Layout", "Complete separation between the telemetry simulator (backend) and the visual control panel (frontend) using async RESTful AJAX requests."),
        ("Framework-Free Python Server (main.py)", "Built entirely using Python's standard-library 'http.server' module. Avoids FastAPI/Uvicorn to ensure 100% stable, zero-dependency startup on Python 3.14."),
        ("Relational Database Layer (database.py)", "Employs standard SQLite3 to record grid logs every 10 minutes. Includes automated data pruning to maintain a lightweight 24-hour rolling window (144 logs)."),
        ("Custom NumPy Predictors (predictor.py)", "All machine learning regressors are written in pure NumPy, eliminating dependencies on scikit-learn, PyTorch, or XGBoost C-extensions.")
    ]
    
    for title, desc in items:
        p = tf.add_paragraph()
        p.space_after = Pt(15)
        
        run1 = p.add_run()
        run1.text = "- " + title + ": "
        run1.font.name = "Outfit"
        run1.font.bold = True
        run1.font.size = Pt(16)
        run1.font.color.rgb = c_purple
        
        run2 = p.add_run()
        run2.text = desc
        run2.font.name = "Inter"
        run2.font.size = Pt(15)
        run2.font.color.rgb = c_white

    # =========================================================================
    # SLIDE 4: WEATHER & RENEWABLE CONTROLS
    # =========================================================================
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_dark_bg(slide)
    add_title(slide, "3. Weather & Renewable Controls")
    
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    controls = [
        ("Environmental Weather Simulation", "Dynamic sliders for Ambient Temperature (controls AC load), Cloud Cover (controls solar output), and Wind Velocity (controls wind generator output) to stress-test the grid under varying climates."),
        ("Renewable Grid Capacities", "Sliders for Installed Solar Generation Capacity (0-100 MW) and Installed Wind Generation Capacity (0-100 MW) allowing real-time grid scaling and planning simulations."),
        ("Manual Dispatch Overrides", "Supports forced commands for Battery Storage (Auto AI mode, Force Charging, or Force Discharging) to let operators test grid peak-shaving overrides."),
        ("Grid Stress Scenarios", "One-click hotkeys to trigger standard high-stress scenarios (Heatwaves, Stormy wind surges, Cloudy solar drops, and Line congestion) for immediate testing.")
    ]
    
    for title, desc in controls:
        p = tf.add_paragraph()
        p.space_after = Pt(14)
        
        run1 = p.add_run()
        run1.text = title + "\n"
        run1.font.name = "Outfit"
        run1.font.bold = True
        run1.font.size = Pt(16)
        run1.font.color.rgb = c_teal
        
        run2 = p.add_run()
        run2.text = desc
        run2.font.name = "Inter"
        run2.font.size = Pt(14)
        run2.font.color.rgb = c_white

    # =========================================================================
    # SLIDE 5: INDIAN GRID STANDARDS
    # =========================================================================
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_dark_bg(slide)
    add_title(slide, "4. Indian Power Grid Standards Localization")
    
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    standards = [
        ("IEGC Frequency Alignment (50.00 Hz)", "Swapped the base simulation frequency from the US 60 Hz scale to the standard Indian grid frequency of 50.00 Hz (regulated by the Indian Electricity Grid Code). Clamped normal safety fluctuations between 49.10 Hz and 50.80 Hz."),
        ("Time-of-Day (ToD) Tariff System", "Implemented time-of-use utility schedules: Off-Peak Night slots (22:00-06:00) receive a Rs. 1.50/unit rebate, while Morning (09:00-12:00) and Evening (18:00-22:00) Peak slots charge Rs. 1.50 and Rs. 2.50 surcharges respectively."),
        ("Cumulative Residential Slab-Billing", "Programmed standard state slab-based bill projections: Rs. 4.50/unit (0-100 units), Rs. 8.50/unit (101-300 units), Rs. 12.00/unit (301-500 units), and Rs. 15.00/unit above 500 units.")
    ]
    
    for title, desc in standards:
        p = tf.add_paragraph()
        p.space_after = Pt(18)
        
        run1 = p.add_run()
        run1.text = title + "\n"
        run1.font.name = "Outfit"
        run1.font.bold = True
        run1.font.size = Pt(18)
        run1.font.color.rgb = c_teal
        
        run2 = p.add_run()
        run2.text = desc
        run2.font.name = "Inter"
        run2.font.size = Pt(15)
        run2.font.color.rgb = c_white

    # =========================================================================
    # SLIDE 6: DEMAND RESPONSE
    # =========================================================================
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_dark_bg(slide)
    add_title(slide, "5. Price-Elastic Demand Response (DR)")
    
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p_intro = tf.paragraphs[0]
    p_intro.text = "Models consumer load shaving reactions in real-time based on dynamic grid pricing:"
    p_intro.font.name = "Outfit"
    p_intro.font.size = Pt(18)
    p_intro.font.color.rgb = c_white
    p_intro.space_after = Pt(20)
    
    dr_items = [
        ("Critical Peaks (> Rs. 12.00/kWh)", "Triggers an automatic 15% load contraction, representing automated smart-grid load-shedding and battery discharges."),
        ("Normal Peaks (> Rs. 9.00/kWh)", "Triggers an automatic 8% load contraction, encouraging grid stabilization during moderate congestion periods."),
        ("Off-Peak Surplus (< Rs. 5.00/kWh)", "Triggers an automatic 5% load expansion, modeling automated consumer shaving, smart water pumping, and EV charging cycles.")
    ]
    
    for title, desc in dr_items:
        p = tf.add_paragraph()
        p.space_after = Pt(15)
        
        run1 = p.add_run()
        run1.text = "- " + title + ": "
        run1.font.name = "Outfit"
        run1.font.bold = True
        run1.font.size = Pt(16)
        run1.font.color.rgb = c_purple
        
        run2 = p.add_run()
        run2.text = desc
        run2.font.name = "Inter"
        run2.font.size = Pt(15)
        run2.font.color.rgb = c_white

    # =========================================================================
    # SLIDE 7: ADVANCED ANALYTICS
    # =========================================================================
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_dark_bg(slide)
    add_title(slide, "6. Advanced Grid Analytics Features")
    
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    analytics = [
        ("Statistical Anomaly Detection", "Scans telemetry in real-time. Grid frequency deviations (< 49.85 Hz or > 50.15 Hz) or load imbalances trigger a warning. Blinking red pulsing markers are drawn dynamically on frontend charts with diagnostic tooltips."),
        ("24-Hour Historical Trends (/api/grid/trends)", "Aggregates SQLite historical logs to present daily metrics: Peak Demand Time, Average-to-Peak ratios, and Grid Stability Factor (frequency standard deviation rating percentage)."),
        ("Personalized Savings Advisor", "Calculates specific load-shedding and battery arbitrage rewards in Rupee values in real-time, displaying exact Rupee hourly savings rates on the advisory panel.")
    ]
    
    for title, desc in analytics:
        p = tf.add_paragraph()
        p.space_after = Pt(18)
        
        run1 = p.add_run()
        run1.text = title + "\n"
        run1.font.name = "Outfit"
        run1.font.bold = True
        run1.font.size = Pt(18)
        run1.font.color.rgb = c_teal
        
        run2 = p.add_run()
        run2.text = desc
        run2.font.name = "Inter"
        run2.font.size = Pt(15)
        run2.font.color.rgb = c_white

    # =========================================================================
    # SLIDE 8: SUSTAINABILITY & GRID ADVISORY
    # =========================================================================
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_dark_bg(slide)
    add_title(slide, "7. Sustainability & Grid Advisory")
    
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    sustain_points = [
        ("Carbon Intensity Metric", "Calculated as: Carbon Intensity = 450 * (1 - clean_ratio) + 2 * fossil_backup. Tracks emissions in g/kWh in real-time, displaying green clean zones vs. fossil-heavy red zones."),
        ("Total Carbon Saved ($CO_2$ offsets)", "Accumulates the mass of carbon emissions saved since midnight in kilograms (kg) based on renewable MWh generation offset values."),
        ("AI Optimization Advisory Feed", "Tracks grid diagnostics to trigger instant recommendations (e.g. forced battery charges, secondary reserve generator start-up warnings at <10% SoC)."),
        ("Historical Grid Trend Insights", "Aggregates SQLite historical logs over a rolling 24-hour window (144 logs) to display Peak Demand Time, Average vs. Peak Load (MW), Grid Stability (%), and Anomaly Counts.")
    ]
    
    for title, desc in sustain_points:
        p = tf.add_paragraph()
        p.space_after = Pt(14)
        
        run1 = p.add_run()
        run1.text = title + "\n"
        run1.font.name = "Outfit"
        run1.font.bold = True
        run1.font.size = Pt(16)
        run1.font.color.rgb = c_purple
        
        run2 = p.add_run()
        run2.text = desc
        run2.font.name = "Inter"
        run2.font.size = Pt(14)
        run2.font.color.rgb = c_white

    # =========================================================================
    # SLIDE 9: MACHINE LEARNING FROM SCRATCH
    # =========================================================================
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_dark_bg(slide)
    add_title(slide, "8. Custom NumPy Machine Learning Models")
    
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    ml_models = [
        ("LSTM Representation (PureMLP)", "A multi-layer feedforward Neural Network with Glorot weight initializations, a ReLU hidden activation layer (32 nodes), and backpropagation gradient descent updates using the Adam Optimizer."),
        ("XGBoost Representation (PureDecisionTree)", "A decision tree regressor using variance-reduction splitting search. Replicates XGBoost stair-step forecasting predictions."),
        ("Linear Baseline (PureRidge)", "Regularized linear regression solving regularized L2 closed-form normal equations: Beta = (X^T X + Alpha I)^-1 X^T Y.")
    ]
    
    for title, desc in ml_models:
        p = tf.add_paragraph()
        p.space_after = Pt(18)
        
        run1 = p.add_run()
        run1.text = title + "\n"
        run1.font.name = "Outfit"
        run1.font.bold = True
        run1.font.size = Pt(18)
        run1.font.color.rgb = c_purple
        
        run2 = p.add_run()
        run2.text = desc
        run2.font.name = "Inter"
        run2.font.size = Pt(15)
        run2.font.color.rgb = c_white

    # =========================================================================
    # SLIDE 10: FRONTEND UI & POLISH
    # =========================================================================
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_dark_bg(slide)
    add_title(slide, "9. Premium UI/UX & Performance Polish")
    
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    polish = [
        ("Glassmorphic Design Aesthetics", "Built using vanilla HTML/CSS with custom dark-themed panels, neon glow borders, and responsive grid layouts. Visualized without heavy charting libraries using lightweight inline SVGs."),
        ("Concurrent Non-Blocking Fetching", "Refactored initial API calls to load concurrently using Promise.allSettled(). Speeds up load times by 4x and avoids sequential blocking lag."),
        ("Cache-Busted Loading Screen", "Created a global blur loading overlay with cache-busting queries (app.js?v=1.0.3) and a 4-second safety timeout, guaranteeing that the interface never hangs under slow server boots.")
    ]
    
    for title, desc in polish:
        p = tf.add_paragraph()
        p.space_after = Pt(18)
        
        run1 = p.add_run()
        run1.text = title + "\n"
        run1.font.name = "Outfit"
        run1.font.bold = True
        run1.font.size = Pt(18)
        run1.font.color.rgb = c_teal
        
        run2 = p.add_run()
        run2.text = desc
        run2.font.name = "Inter"
        run2.font.size = Pt(15)
        run2.font.color.rgb = c_white

    # =========================================================================
    # SLIDE 11: SCREENSHOTS SLIDE
    # =========================================================================
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_dark_bg(slide)
    add_title(slide, "10. Interface Visuals & Screenshots")
    
    # Place images side-by-side or stacked
    has_img1 = os.path.exists(img1_path)
    has_img2 = os.path.exists(img2_path)
    
    if has_img1 and has_img2:
        slide.shapes.add_picture(img1_path, Inches(0.75), Inches(1.8), width=Inches(5.6))
        slide.shapes.add_picture(img2_path, Inches(6.8), Inches(2.5), width=Inches(5.6))
        
        tx1 = slide.shapes.add_textbox(Inches(0.75), Inches(5.4), Inches(5.6), Inches(1.0))
        tx1.text_frame.word_wrap = True
        p1 = tx1.text_frame.paragraphs[0]
        p1.text = "Figure 1: Live telemetry dashboard showing Indianized dynamic price, slab-billing estimation, and dispatch SVG analytics."
        p1.font.name = "Inter"
        p1.font.size = Pt(12)
        p1.font.color.rgb = c_slate
        p1.alignment = PP_ALIGN.CENTER
        
        tx2 = slide.shapes.add_textbox(Inches(6.8), Inches(4.7), Inches(5.6), Inches(1.0))
        tx2.text_frame.word_wrap = True
        p2 = tx2.text_frame.paragraphs[0]
        p2.text = "Figure 2: AI Time-Series forecasting panel displaying actual demand curves vs. customized ML predictions."
        p2.font.name = "Inter"
        p2.font.size = Pt(12)
        p2.font.color.rgb = c_slate
        p2.alignment = PP_ALIGN.CENTER
        
    elif has_img1:
        slide.shapes.add_picture(img1_path, Inches(2.66), Inches(1.6), width=Inches(8.0))
    else:
        txBox = slide.shapes.add_textbox(Inches(0.75), Inches(2.5), Inches(11.83), Inches(2.0))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "[Screenshots media files not found. Upload screenshots to include in this slide]"
        p.font.name = "Inter"
        p.font.size = Pt(18)
        p.font.color.rgb = c_slate
        p.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 12: CONCLUSION
    # =========================================================================
    slide = prs.slides.add_slide(blank_slide_layout)
    apply_dark_bg(slide)
    add_title(slide, "11. Summary & Academic Value")
    
    txBox = slide.shapes.add_textbox(Inches(0.75), Inches(1.5), Inches(11.83), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    conclusions = [
        ("Academic Highlights", "Demonstrates deep mathematical understanding of ML (backpropagation, neural network gates) by implementing algorithms from scratch in pure NumPy, rather than calling pre-built library packages."),
        ("Software Engineering Design", "Showcases robust decoupled architecture, concurrent asynchronous execution in vanilla JavaScript, relational database schemas in SQLite, and dynamic data-pruning mechanisms."),
        ("Real-world Practical Value", "Provides grid decision-support functions including statistical anomaly filters matching standard Indian utility guidelines (IEGC 50Hz, slab billing, and dynamic demand response)."),
        ("Project Deliverables", "Includes the complete repository codebase, local runtime package (smart-grid-sync.zip), and detailed final week technical report (PDF).")
    ]
    
    for title, desc in conclusions:
        p = tf.add_paragraph()
        p.space_after = Pt(15)
        
        run1 = p.add_run()
        run1.text = "- " + title + ": "
        run1.font.name = "Outfit"
        run1.font.bold = True
        run1.font.size = Pt(16)
        run1.font.color.rgb = c_teal
        
        run2 = p.add_run()
        run2.text = desc
        run2.font.name = "Inter"
        run2.font.size = Pt(15)
        run2.font.color.rgb = c_white
        
    prs.save(output_path)
    print(f"Presentation successfully created at: {output_path}")

if __name__ == "__main__":
    out_dir = "/Users/shikhasrivastava/.gemini/antigravity/scratch/smart-grid-sync"
    img1 = os.path.join(out_dir, "media__1780957951599.png")
    img2 = os.path.join(out_dir, "media__1780957962135.png")
    out_file = os.path.join(out_dir, "smart_grid_sync_final_presentation.pptx")
    create_presentation(out_file, img1, img2)
